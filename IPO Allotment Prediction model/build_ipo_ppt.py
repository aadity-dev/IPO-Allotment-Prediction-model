from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------- CONFIG ----------------
TITLE = "IPO Allotment Prediction"
SUBTITLE = "Using Machine Learning (2007–2025)"
AUTHOR = "ADITYA (12318430)"
FACULTY = "Submitted to: Mahipal Sir"

# Palette
THEME_PRIMARY = RGBColor(0, 123, 167)   # teal-blue
THEME_DARK = RGBColor(13, 17, 23)       # near-black
THEME_MUTED = RGBColor(100, 116, 139)   # gray

OUTPUT = "IPO-Allotment-Prediction.pptx"

import os

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))

# Image filenames (ensure they exist; change here if your names differ)
IMG_CM_OVERALL = os.path.join(SCRIPT_DIR, "report_assets", "confusion_matrix.png")
IMG_CM_OVERALL_2 = os.path.join(SCRIPT_DIR, "report_assets", "confusion matrix.png")
IMG_CM_LR = os.path.join(SCRIPT_DIR, "report_assets", "Logistic_Regression_confusion_matrix.png")
IMG_CM_RF = os.path.join(SCRIPT_DIR, "report_assets", "Random_Forest_confusion_matrix.png")
IMG_ROC_LR = os.path.join(SCRIPT_DIR, "report_assets", "Logistic_Regression_roc_curve.png")
IMG_ROC_RF = os.path.join(SCRIPT_DIR, "report_assets", "Random_Forest_roc_curve.png")
IMG_PAIRPLOT = os.path.join(SCRIPT_DIR, "report_assets", "pairplot.png")
IMG_BOX_SUB = os.path.join(SCRIPT_DIR, "report_assets", "subscription_by_category_boxplot.png")
IMG_TARGET_DIST = os.path.join(SCRIPT_DIR, "report_assets", "target_distribution.png")

# ---------------- HELPERS ----------------
def add_footer(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(6.85), Inches(8.6), Inches(0.4))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = THEME_MUTED

def style_title_shape(title_shape, size_pt=34):
    if title_shape and title_shape.has_text_frame:
        tf = title_shape.text_frame
        if tf.paragraphs:
            f = tf.paragraphs[0].font
            f.size = Pt(size_pt)
            f.bold = True
            f.color.rgb = THEME_DARK


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    title_shape = slide.shapes.title
    if title_shape is None:
        title_shape = slide.shapes.add_textbox(Inches(0.9), Inches(0.7), Inches(8.2), Inches(1.0))
    title_shape.text = title_text
    style_title_shape(title_shape, 40)

    # Subtitle
    sub = slide.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(8.2), Inches(1.0))
    tf = sub.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = subtitle_text
    p.font.size = Pt(20)
    p.font.color.rgb = THEME_MUTED

    # Author / Faculty
    meta = slide.shapes.add_textbox(Inches(0.9), Inches(2.6), Inches(8.2), Inches(1.2))
    tmeta = meta.text_frame
    tmeta.clear()
    p1 = tmeta.paragraphs[0]
    p1.text = AUTHOR
    p1.font.size = Pt(18); p1.font.color.rgb = THEME_DARK
    p2 = tmeta.add_paragraph()
    p2.text = FACULTY
    p2.font.size = Pt(16); p2.font.color.rgb = THEME_MUTED

    add_footer(slide, "Academic Project | Data Mining & ML")
    slide.notes_slide.notes_text_frame.text = "Introduce scope and that Random Forest achieved ROC-AUC 0.83."
    return slide

def add_bullets_slide(prs, title, bullets, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide.shapes.title.text = title
    style_title_shape(slide.shapes.title, 34)
    ph = slide.shapes.placeholders[1]
    tf = ph.text_frame
    tf.clear()
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(20)
    add_footer(slide, f"{AUTHOR}  |  {FACULTY}")
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide

def add_image_slide(prs, title, img_path, notes=None, width=Inches(9.2)):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide.shapes.title.text = title
    style_title_shape(slide.shapes.title, 34)
    left = (Inches(10) - width) / 2
    top = Inches(1.6)
    slide.shapes.add_picture(img_path, left, top, width=width)
    add_footer(slide, f"{AUTHOR}  |  {FACULTY}")
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide

def add_side_by_side_images(prs, title, img_left, img_right, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    style_title_shape(slide.shapes.title, 34)
    w = Inches(4.4)
    left1 = Inches(0.6)
    left2 = Inches(5.0)
    top = Inches(1.6)
    slide.shapes.add_picture(img_left, left1, top, width=w)
    slide.shapes.add_picture(img_right, left2, top, width=w)
    add_footer(slide, f"{AUTHOR}  |  {FACULTY}")
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide

# ---------------- BUILD ----------------
def build():
    prs = Presentation()

    # 1) Title
    add_title_slide(prs, TITLE, SUBTITLE)

    # 2) Motivation & Objective
    add_bullets_slide(
        prs, "Motivation & Objective",
        [
            "IPO allotments are competitive and uncertain.",
            "Provide data-driven guidance on allotment success likelihood.",
            "Binary classification: predict success (1) vs failure (0)."
        ],
        "Stress that predictions are probabilistic, not guarantees."
    )

    # 3) Data Overview
    add_bullets_slide(
        prs, "Data Overview (2007–2025)",
        [
            "Merged two CSVs via fuzzy company-name matching.",
            "Target: Allotment outcome (1/0).",
            "Key signals: subscription rates, issue size, listing gains, investor category."
        ],
        "Keep it high-level; details later."
    )

    # 4) Preprocessing
    add_bullets_slide(
        prs, "Preprocessing Pipeline",
        [
            "Fuzzy merge on names; handle missing/standardize columns.",
            "Feature engineering: Log(Issue Size), Investor Category bins.",
            "One-Hot encoding; standardize numerical features.",
            "Train/Test split: 80/20."
        ],
        "Mention scaling and encoding choices briefly."
    )

    # 5) Key Features
    add_bullets_slide(
        prs, "Key Features",
        [
            "Issue Size (log-transformed).",
            "Subscription rates: QIB, HNI, RII, Total.",
            "Listing Gain (%).",
            "Investor Category (binned)."
        ],
        "Subscription intensity and size drive predictions."
    )

    # 6) EDA: Target Distribution
    add_image_slide(prs, "EDA: Target Distribution", IMG_TARGET_DIST,
                    "Class balance overview; note any imbalance.")

    # 7) EDA: Pairplot
    add_image_slide(prs, "EDA: Key Signals vs Target (Pairplot)", IMG_PAIRPLOT,
                    "Point out separability patterns.")

    # 8) EDA: Subscription by Category
    add_image_slide(prs, "EDA: Subscription by Investor Category (Boxplot)", IMG_BOX_SUB,
                    "HNI vs Retail intensity differences.")

    # 9) Modeling Approach
    add_bullets_slide(
        prs, "Modeling Approach",
        [
            "Logistic Regression (baseline linear).",
            "Random Forest (ensemble of decision trees).",
            "XGBoost (gradient boosting)."
        ],
        "Focus on intuition over math."
    )

    # 10) Evaluation Protocol
    add_bullets_slide(
        prs, "Evaluation Protocol",
        [
            "Split: 80% train, 20% test.",
            "Primary metric: ROC-AUC.",
            "Also: Confusion Matrix, Precision/Recall, F1."
        ],
        "Why ROC-AUC: threshold-agnostic discrimination."
    )

    # 11) Results Summary (AUC)
    add_bullets_slide(
        prs, "Results Summary (Test ROC-AUC)",
        [
            "Logistic Regression: 0.51",
            "Random Forest: 0.83",
            "XGBoost: 0.63"
        ],
        "Random Forest performed best."
    )

    # 12) ROC Curves
    add_side_by_side_images(prs, "ROC Curves: Logistic Regression vs Random Forest",
                            IMG_ROC_LR, IMG_ROC_RF,
                            "Contrast separability (AUC).")

    # 13) Confusion Matrices by Model
    add_side_by_side_images(prs, "Confusion Matrices: Logistic Regression vs Random Forest",
                            IMG_CM_LR, IMG_CM_RF,
                            "Discuss false positives/negatives.")

    # 14) Overall Confusion Matrices (optional views)
    add_side_by_side_images(prs, "Additional Confusion Matrices (Alternative Views)",
                            IMG_CM_OVERALL, IMG_CM_OVERALL_2,
                            "Supplementary displays if relevant.")

    # 15) Discussion & Investor Value
    add_bullets_slide(
        prs, "Discussion & Investor Value",
        [
            "Subscription intensity and issue size drive predictions.",
            "Ensembles capture nonlinearity better than linear baseline.",
            "Use probabilities to prioritize IPO applications."
        ],
        "Be transparent about inherent randomness."
    )

    # 16) Conclusion & Future Work
    add_bullets_slide(
        prs, "Conclusion & Future Work",
        [
            "Random Forest is most effective (ROC-AUC 0.83).",
            "Add sentiment and underwriter data; calibrate probabilities.",
            "Consider SHAP for explainability; time-based validation."
        ],
        "End with actionable roadmap."
    )

    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")

if __name__ == "__main__":
    build()
